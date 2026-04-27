import importlib.util

import pandas as pd
import pytest
from docx import Document
from PIL import Image
from pptx import Presentation
from pypdf import PdfWriter

from app.model.scraper.document_scraper import DocumentRow, DocumentScraper


HAS_EBOOKLIB = importlib.util.find_spec("ebooklib") is not None
HAS_OPENPYXL = importlib.util.find_spec("openpyxl") is not None
HAS_STRIPRTF = importlib.util.find_spec("striprtf") is not None


class FakeHeaders:
    def get_content_charset(self):
        return "utf-8"


class FakeUrlResponse:
    headers = FakeHeaders()

    def __init__(self, body):
        self.body = body

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, traceback):
        return False

    def read(self):
        return self.body


@pytest.fixture
def scraper():
    return DocumentScraper()


def build_pdf_bytes(text):
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 144] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        (
            b"<< /Length %d >>\nstream\nBT\n/F1 18 Tf\n50 100 Td\n(%s) Tj\nET\nendstream"
            % (
                len(
                    f"BT\n/F1 18 Tf\n50 100 Td\n({text}) Tj\nET\n".encode(
                        "utf-8"
                    )
                ),
                text.encode("utf-8"),
            )
        ),
    ]

    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{index} 0 obj\n".encode("utf-8"))
        pdf.extend(obj)
        pdf.extend(b"\nendobj\n")

    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {len(objects) + 1}\n".encode("utf-8"))
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("utf-8"))
    pdf.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("utf-8")
    )
    return bytes(pdf)


def test_scrape_txt_splits_into_document_rows(runtime_tmp_path, scraper):
    text_path = runtime_tmp_path / "sample.txt"
    text_path.write_text("First block.\n\nSecond block.", encoding="utf-8")

    rows = scraper.scrape_file(text_path)

    assert len(rows) == 2
    assert rows[0]["title"] == "Text Block 1"
    assert rows[0]["primary_text"] == "First block."
    assert rows[1]["primary_text"] == "Second block."
    assert rows[0]["source_type"] == "txt"
    assert rows[0]["metadata"] == {}


def test_dependency_status_reports_missing_advertised_parser_packages(monkeypatch):
    missing_modules = {"ebooklib", "striprtf", "xlrd"}

    def fake_find_spec(module_name):
        return None if module_name in missing_modules else object()

    monkeypatch.setattr(
        "app.model.scraper.document_scraper.importlib.util.find_spec",
        fake_find_spec,
    )

    missing_dependencies = DocumentScraper.missing_runtime_dependencies()

    missing_packages = {dependency.package_name for dependency in missing_dependencies}
    assert {"ebooklib", "striprtf", "xlrd"}.issubset(missing_packages)


def test_scrape_file_fails_clearly_when_parser_dependency_is_missing(
    runtime_tmp_path, scraper, monkeypatch
):
    rtf_path = runtime_tmp_path / "sample.rtf"
    rtf_path.write_text(r"{\rtf1\ansi Missing dependency sample.}", encoding="utf-8")

    def fake_find_spec(module_name):
        return None if module_name == "striprtf" else object()

    monkeypatch.setattr(
        "app.model.scraper.document_scraper.importlib.util.find_spec",
        fake_find_spec,
    )

    with pytest.raises(RuntimeError, match="striprtf"):
        scraper.scrape_file(rtf_path)


def test_xls_imports_require_the_xlrd_reader_package():
    missing_dependencies = DocumentScraper.missing_dependencies_for_suffix(".xls")
    package_names = {dependency.package_name for dependency in missing_dependencies}

    if importlib.util.find_spec("xlrd") is None:
        assert "xlrd" in package_names
    else:
        assert "xlrd" not in package_names


def test_scrape_html_extracts_text_sections(runtime_tmp_path, scraper):
    html_path = runtime_tmp_path / "sample.html"
    html_path.write_text(
        (
            "<html><head><title>Sample</title></head><body>"
            "<h1>Heading</h1><p>Paragraph one.</p><p>Paragraph two.</p>"
            "</body></html>"
        ),
        encoding="utf-8",
    )

    rows = scraper.scrape_file(html_path)

    assert len(rows) == 1
    assert rows[0]["title"] == "Heading"
    assert rows[0]["source_type"] == "html"
    assert "Heading" in rows[0]["primary_text"]
    assert "Paragraph two." in rows[0]["primary_text"]
    assert rows[0]["metadata"]["kind"] == "html_section"
    assert rows[0]["metadata"]["heading_level"] == 1


def test_scrape_html_preserves_table_rows(runtime_tmp_path, scraper):
    html_path = runtime_tmp_path / "table.html"
    html_path.write_text(
        (
            "<html><body>"
            "<table><tr><th>Step</th><th>Status</th></tr>"
            "<tr><td>Import</td><td>Ready</td></tr></table>"
            "</body></html>"
        ),
        encoding="utf-8",
    )

    rows = scraper.scrape_file(html_path)

    assert len(rows) == 1
    assert rows[0]["title"] == "Table 1 Row 1"
    assert "Step: Import" in rows[0]["primary_text"]
    assert rows[0]["metadata"]["kind"] == "table_row"


def test_scrape_html_text_extracts_pasted_html_sections(scraper):
    rows = scraper.scrape_html_text(
        (
            "<html><head><title>Pasted Source</title></head><body>"
            "<h2>Raw HTML Heading</h2><p>Unique pasted HTML body.</p>"
            "</body></html>"
        )
    )

    assert len(rows) == 1
    assert rows[0]["title"] == "Raw HTML Heading"
    assert rows[0]["source_path"] == "raw-html"
    assert rows[0]["source_type"] == "html"
    assert "Unique pasted HTML body." in rows[0]["primary_text"]


def test_scrape_url_fetches_and_scrapes_html(scraper, monkeypatch):
    html = (
        "<html><head><title>Remote Source</title></head><body>"
        "<h1>Remote Heading</h1><p>Unique URL imported body.</p>"
        "</body></html>"
    ).encode("utf-8")
    calls = []

    def fake_urlopen(request, timeout):
        calls.append((request, timeout))
        return FakeUrlResponse(html)

    monkeypatch.setattr("app.model.scraper.document_scraper.urlopen", fake_urlopen)

    rows = scraper.scrape_url("https://example.test/article")

    assert len(calls) == 1
    assert len(rows) == 1
    assert rows[0]["title"] == "Remote Heading"
    assert rows[0]["source_path"] == "https://example.test/article"
    assert rows[0]["source_type"] == "url"
    assert "Unique URL imported body." in rows[0]["primary_text"]


def test_scrape_url_rejects_non_web_schemes(scraper):
    with pytest.raises(ValueError, match="http or https"):
        scraper.scrape_url("file:///tmp/example.html")


def test_scrape_docx_extracts_paragraphs(runtime_tmp_path, scraper):
    docx_path = runtime_tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("First paragraph.")
    document.add_paragraph("Second paragraph.")
    document.save(docx_path)

    rows = scraper.scrape_file(docx_path)

    assert len(rows) == 2
    assert rows[0]["title"] == "Paragraph 1"
    assert rows[1]["primary_text"] == "Second paragraph."
    assert rows[0]["metadata"]["kind"] == "paragraph"


def test_scrape_docx_preserves_headings_and_tables(runtime_tmp_path, scraper):
    docx_path = runtime_tmp_path / "structured.docx"
    document = Document()
    document.add_heading("Launch Plan", level=1)
    document.add_paragraph("Prepare the narration script.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Task"
    table.cell(0, 1).text = "Owner"
    table.cell(1, 0).text = "Record intro"
    table.cell(1, 1).text = "Narrator"
    document.save(docx_path)

    rows = scraper.scrape_file(docx_path)

    assert rows[0]["title"] == "Launch Plan"
    assert "Prepare the narration script." in rows[0]["primary_text"]
    assert rows[0]["metadata"]["kind"] == "section"
    assert rows[0]["metadata"]["heading_level"] == 1
    assert rows[1]["title"] == "Table 1 Row 1"
    assert "Task: Record intro" in rows[1]["primary_text"]
    assert rows[1]["metadata"]["kind"] == "table_row"


def test_scrape_pdf_extracts_page_text(runtime_tmp_path, scraper):
    pdf_path = runtime_tmp_path / "sample.pdf"
    pdf_path.write_bytes(build_pdf_bytes("Hello PDF"))

    rows = scraper.scrape_file(pdf_path)

    assert len(rows) == 1
    assert rows[0]["title"] == "Page 1"
    assert "Hello PDF" in rows[0]["primary_text"]
    assert rows[0]["metadata"]["kind"] == "page"
    assert rows[0]["metadata"]["extraction"] == "text-layer"


@pytest.mark.skipif(not HAS_STRIPRTF, reason="striprtf is not installed")
def test_scrape_rtf_extracts_text(runtime_tmp_path, scraper):
    rtf_path = runtime_tmp_path / "sample.rtf"
    rtf_path.write_text(r"{\rtf1\ansi This is {\b bold} text.}", encoding="utf-8")

    rows = scraper.scrape_file(rtf_path)

    assert len(rows) == 1
    assert rows[0]["title"] == "RTF Block 1"
    assert "This is bold text." in rows[0]["primary_text"]


@pytest.mark.skipif(not HAS_EBOOKLIB, reason="ebooklib is not installed")
def test_scrape_epub_extracts_chapter_text(runtime_tmp_path, scraper):
    from ebooklib import epub

    epub_path = runtime_tmp_path / "sample.epub"
    book = epub.EpubBook()
    book.set_identifier("sample-book")
    book.set_title("Sample Book")
    book.set_language("en")
    chapter = epub.EpubHtml(title="Chapter 1", file_name="chapter1.xhtml", lang="en")
    chapter.content = (
        "<html><head><title>Chapter 1</title></head>"
        "<body><p>EPUB body text.</p></body></html>"
    )
    book.add_item(chapter)
    book.toc = (chapter,)
    book.spine = ["nav", chapter]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(epub_path), book)

    rows = scraper.scrape_file(epub_path)

    assert len(rows) == 1
    assert rows[0]["title"] == "Chapter 1"
    assert "EPUB body text." in rows[0]["primary_text"]
    assert rows[0]["metadata"]["kind"] == "chapter"
    assert rows[0]["metadata"]["chapter_number"] == 1


def test_scrape_csv_extracts_rows(runtime_tmp_path, scraper):
    csv_path = runtime_tmp_path / "sample.csv"
    pd.DataFrame(
        [{"Name": "Alice", "Role": "Narrator"}, {"Name": "Bob", "Role": "Editor"}]
    ).to_csv(csv_path, index=False)

    rows = scraper.scrape_file(csv_path)

    assert len(rows) == 2
    assert rows[0]["title"] == "Sheet1 Row 1"
    assert "Name: Alice" in rows[0]["primary_text"]
    assert "Columns: Name, Role" in rows[0]["secondary_text"]
    assert rows[0]["metadata"]["kind"] == "sheet_row"
    assert rows[0]["metadata"]["sheet_name"] == "Sheet1"


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl is not installed")
def test_scrape_xlsx_extracts_rows(runtime_tmp_path, scraper):
    xlsx_path = runtime_tmp_path / "sample.xlsx"
    with pd.ExcelWriter(xlsx_path) as writer:
        pd.DataFrame([{"Chapter": "Intro", "Pages": 2}]).to_excel(
            writer,
            sheet_name="Outline",
            index=False,
        )

    rows = scraper.scrape_file(xlsx_path)

    assert len(rows) == 1
    assert rows[0]["title"] == "Outline Row 1"
    assert "Chapter: Intro" in rows[0]["primary_text"]
    assert rows[0]["metadata"]["columns"] == ["Chapter", "Pages"]


def test_scrape_pptx_extracts_slide_text_and_notes(runtime_tmp_path, scraper):
    pptx_path = runtime_tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide Heading"
    slide.placeholders[1].text = "Slide body text."
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "Speaker notes."
    presentation.save(pptx_path)

    rows = scraper.scrape_file(pptx_path)

    assert len(rows) == 1
    assert rows[0]["title"] == "Slide 1"
    assert "Slide Heading" in rows[0]["primary_text"]
    assert "Speaker notes." in rows[0]["secondary_text"]
    assert rows[0]["metadata"]["kind"] == "slide"
    assert rows[0]["metadata"]["has_notes"] is True


def test_scrape_image_uses_ocr_path(runtime_tmp_path, scraper, monkeypatch):
    image_path = runtime_tmp_path / "sample.png"
    Image.new("RGB", (40, 20), color="white").save(image_path)
    monkeypatch.setattr(
        "app.model.scraper.document_scraper.importlib.util.find_spec",
        lambda _module_name: object(),
    )
    ocr_calls = []

    def fake_ocr(self, path):
        ocr_calls.append(path)
        return "Scanned image text"

    monkeypatch.setattr(DocumentScraper, "_perform_ocr_on_image", fake_ocr)

    rows = scraper.scrape_file(image_path)

    assert ocr_calls == [image_path]
    assert len(rows) == 1
    assert rows[0]["title"] == "sample"
    assert rows[0]["primary_text"] == "Scanned image text"
    assert rows[0]["source_type"] == "png"
    assert rows[0]["metadata"]["extraction"] == "ocr"


def test_scrape_pdf_falls_back_to_ocr_when_text_layer_is_missing(
    runtime_tmp_path, scraper, monkeypatch
):
    pdf_path = runtime_tmp_path / "scanned.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with open(pdf_path, "wb") as file:
        writer.write(file)
    ocr_calls = []

    def fake_pdf_ocr(self, path):
        ocr_calls.append(path)
        return [DocumentRow(1, "Page 1", "OCR page text")]

    monkeypatch.setattr(DocumentScraper, "_perform_ocr_on_pdf", fake_pdf_ocr)

    rows = scraper.scrape_file(pdf_path)

    assert ocr_calls == [pdf_path]
    assert len(rows) == 1
    assert rows[0]["title"] == "Page 1"
    assert rows[0]["primary_text"] == "OCR page text"
    assert rows[0]["metadata"]["extraction"] == "ocr"
