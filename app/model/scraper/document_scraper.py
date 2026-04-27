"""Unified document scraping helpers for the import dialog."""

from __future__ import annotations

from dataclasses import dataclass, field
import importlib.util
from pathlib import Path
from urllib.parse import urlparse
from urllib.request import Request, urlopen


@dataclass
class DocumentRow:
    item_number: int
    title: str
    primary_text: str
    secondary_text: str = ""
    metadata: dict = field(default_factory=dict)


@dataclass(frozen=True)
class ParserDependency:
    module_name: str
    package_name: str
    purpose: str


class DocumentScraper:
    """Scrape common document formats into normalized document rows."""

    SUPPORTED_FILTER = (
        "Supported Documents (*.txt *.docx *.pdf *.html *.htm *.rtf *.epub "
        "*.xlsx *.xls *.csv *.pptx *.png *.jpg *.jpeg *.tif *.tiff *.bmp *.webp);;"
        "Text Documents (*.txt);;"
        "Word Documents (*.docx);;"
        "PDF Files (*.pdf);;"
        "HTML Files (*.html *.htm);;"
        "RTF Files (*.rtf);;"
        "EPUB Files (*.epub);;"
        "Spreadsheet Files (*.xlsx *.xls *.csv);;"
        "PowerPoint Files (*.pptx);;"
        "Image Files (*.png *.jpg *.jpeg *.tif *.tiff *.bmp *.webp);;"
        "All Files (*.*)"
    )
    FORMAT_DEPENDENCIES = {
        ".txt": (),
        ".docx": (
            ParserDependency("docx", "python-docx", "Word document imports"),
        ),
        ".pdf": (
            ParserDependency("pypdf", "pypdf", "PDF text-layer imports"),
        ),
        ".html": (
            ParserDependency("bs4", "beautifulsoup4", "HTML imports"),
        ),
        ".htm": (
            ParserDependency("bs4", "beautifulsoup4", "HTML imports"),
        ),
        ".rtf": (
            ParserDependency("striprtf", "striprtf", "RTF imports"),
        ),
        ".epub": (
            ParserDependency("bs4", "beautifulsoup4", "EPUB HTML extraction"),
            ParserDependency("ebooklib", "ebooklib", "EPUB imports"),
        ),
        ".xlsx": (
            ParserDependency("pandas", "pandas", "spreadsheet imports"),
            ParserDependency("openpyxl", "openpyxl", "XLSX imports"),
        ),
        ".xls": (
            ParserDependency("pandas", "pandas", "spreadsheet imports"),
            ParserDependency("xlrd", "xlrd", "legacy XLS imports"),
        ),
        ".csv": (
            ParserDependency("pandas", "pandas", "CSV imports"),
        ),
        ".pptx": (
            ParserDependency("pptx", "python-pptx", "PowerPoint imports"),
        ),
        ".png": (
            ParserDependency("PIL", "pillow", "image imports"),
            ParserDependency("pytesseract", "pytesseract", "image OCR imports"),
        ),
        ".jpg": (
            ParserDependency("PIL", "pillow", "image imports"),
            ParserDependency("pytesseract", "pytesseract", "image OCR imports"),
        ),
        ".jpeg": (
            ParserDependency("PIL", "pillow", "image imports"),
            ParserDependency("pytesseract", "pytesseract", "image OCR imports"),
        ),
        ".tif": (
            ParserDependency("PIL", "pillow", "image imports"),
            ParserDependency("pytesseract", "pytesseract", "image OCR imports"),
        ),
        ".tiff": (
            ParserDependency("PIL", "pillow", "image imports"),
            ParserDependency("pytesseract", "pytesseract", "image OCR imports"),
        ),
        ".bmp": (
            ParserDependency("PIL", "pillow", "image imports"),
            ParserDependency("pytesseract", "pytesseract", "image OCR imports"),
        ),
        ".webp": (
            ParserDependency("PIL", "pillow", "image imports"),
            ParserDependency("pytesseract", "pytesseract", "image OCR imports"),
        ),
    }
    OCR_PDF_DEPENDENCIES = (
        ParserDependency("pypdfium2", "pypdfium2", "scanned PDF rendering"),
        ParserDependency("PIL", "pillow", "scanned PDF image handling"),
        ParserDependency("pytesseract", "pytesseract", "scanned PDF OCR"),
    )

    @classmethod
    def missing_dependencies_for_suffix(cls, suffix):
        """Return parser dependencies missing for a specific file suffix."""
        normalized_suffix = suffix.lower()
        return [
            dependency
            for dependency in cls.FORMAT_DEPENDENCIES.get(normalized_suffix, ())
            if importlib.util.find_spec(dependency.module_name) is None
        ]

    @classmethod
    def missing_runtime_dependencies(cls):
        """Return missing parser dependencies for advertised import formats."""
        missing_by_package = {}
        for dependencies in cls.FORMAT_DEPENDENCIES.values():
            for dependency in dependencies:
                if importlib.util.find_spec(dependency.module_name) is None:
                    missing_by_package[dependency.package_name] = dependency

        for dependency in cls.OCR_PDF_DEPENDENCIES:
            if importlib.util.find_spec(dependency.module_name) is None:
                missing_by_package[dependency.package_name] = dependency

        return sorted(missing_by_package.values(), key=lambda item: item.package_name)

    @classmethod
    def format_missing_dependency_message(cls, dependencies):
        packages = ", ".join(dependency.package_name for dependency in dependencies)
        return (
            f"Missing parser package(s): {packages}. "
            "Run `poetry install` from the project directory so all advertised "
            "document formats are available."
        )

    @classmethod
    def dependency_status_message(cls):
        missing_dependencies = cls.missing_runtime_dependencies()
        if not missing_dependencies:
            return ""
        return cls.format_missing_dependency_message(missing_dependencies)

    def scrape_file(self, file_path):
        path = Path(file_path)
        suffix = path.suffix.lower()

        if suffix not in self.FORMAT_DEPENDENCIES:
            raise ValueError(f"Unsupported document type: {suffix or '[none]'}")

        missing_dependencies = self.missing_dependencies_for_suffix(suffix)
        if missing_dependencies:
            raise RuntimeError(
                self.format_missing_dependency_message(missing_dependencies)
            )

        if suffix == ".txt":
            rows = self._scrape_txt(path)
        elif suffix == ".docx":
            rows = self._scrape_docx(path)
        elif suffix == ".pdf":
            rows = self._scrape_pdf(path)
        elif suffix in {".html", ".htm"}:
            rows = self._scrape_html(path)
        elif suffix == ".rtf":
            rows = self._scrape_rtf(path)
        elif suffix == ".epub":
            rows = self._scrape_epub(path)
        elif suffix in {".xlsx", ".xls", ".csv"}:
            rows = self._scrape_spreadsheet(path)
        elif suffix == ".pptx":
            rows = self._scrape_pptx(path)
        elif suffix in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}:
            rows = self._scrape_image(path)
        return self._rows_to_dicts(rows, path, suffix.lstrip("."))

    def scrape_html_text(
        self,
        html,
        source_name="Pasted HTML",
        source_path="raw-html",
        source_type="html",
    ):
        """Scrape raw HTML content into the same normalized rows as HTML files."""
        missing_dependencies = self.missing_dependencies_for_suffix(".html")
        if missing_dependencies:
            raise RuntimeError(
                self.format_missing_dependency_message(missing_dependencies)
            )

        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html, "html.parser")
        title = soup.title.get_text(strip=True) if soup.title else source_name
        rows = self._html_to_rows(soup, title, "html_section")
        return self._rows_to_dicts(rows, source_path, source_type)

    def scrape_url(self, url, timeout=20):
        """Fetch a URL and scrape the response body as structured HTML content."""
        parsed_url = urlparse(url)
        if parsed_url.scheme not in {"http", "https"}:
            raise ValueError("URL imports require an http or https address.")

        request = Request(
            url,
            headers={"User-Agent": "TextToSpeechPython document importer"},
        )
        with urlopen(request, timeout=timeout) as response:
            charset = response.headers.get_content_charset() or "utf-8"
            html = response.read().decode(charset, errors="replace")

        return self.scrape_html_text(
            html,
            source_name=url,
            source_path=url,
            source_type="url",
        )

    def _rows_to_dicts(self, rows, source_path, source_type):
        return [
            self._row_to_dict(row, source_path, source_type)
            for row in rows
            if self._has_content(row)
        ]

    def _row_to_dict(self, row, source_path, source_type):
        return {
            "item_number": row.item_number,
            "title": row.title,
            "primary_text": row.primary_text,
            "secondary_text": row.secondary_text,
            "source_path": str(source_path),
            "source_type": source_type,
            "metadata": dict(row.metadata),
        }

    def _has_content(self, row):
        return bool(
            (row.primary_text or "").strip()
            or (row.secondary_text or "").strip()
        )

    def _split_blocks(self, text):
        normalized = text.replace("\r\n", "\n").replace("\r", "\n")
        blocks = [block.strip() for block in normalized.split("\n\n")]
        return [block for block in blocks if block]

    def _heading_level(self, style_name):
        normalized_style = (style_name or "").strip().lower()
        if not normalized_style.startswith("heading"):
            return None

        level_text = normalized_style.replace("heading", "").strip()
        return int(level_text) if level_text.isdigit() else 1

    def _table_to_rows(self, table_rows, start_item_number, table_index):
        extracted_rows = []
        raw_rows = []
        for row in table_rows:
            if hasattr(row, "cells"):
                raw_rows.append([cell.text.strip() for cell in row.cells])
            else:
                raw_rows.append([str(cell).strip() for cell in row])
        raw_rows = [row for row in raw_rows if any(cell for cell in row)]
        if not raw_rows:
            return extracted_rows

        headers = raw_rows[0]
        data_rows = raw_rows[1:] if len(raw_rows) > 1 else raw_rows
        for row_offset, cells in enumerate(data_rows, start=1):
            if len(headers) == len(cells) and len(raw_rows) > 1:
                values = [
                    f"{header or f'Column {index + 1}'}: {value}"
                    for index, (header, value) in enumerate(zip(headers, cells))
                    if value
                ]
            else:
                values = [
                    f"Column {index + 1}: {value}"
                    for index, value in enumerate(cells)
                    if value
                ]

            if not values:
                continue

            extracted_rows.append(
                DocumentRow(
                    start_item_number + len(extracted_rows),
                    f"Table {table_index} Row {row_offset}",
                    "\n".join(values),
                    f"Table: {table_index}",
                    {
                        "kind": "table_row",
                        "table_number": table_index,
                        "row_number": row_offset,
                        "headers": headers if len(raw_rows) > 1 else [],
                    },
                )
            )
        return extracted_rows

    def _html_to_rows(self, soup, document_title, row_kind):
        rows = []
        current_heading = None
        current_level = None
        current_chunks = []
        table_index = 1

        def flush_section():
            nonlocal current_heading, current_level, current_chunks
            if not current_chunks:
                return

            item_number = len(rows) + 1
            rows.append(
                DocumentRow(
                    item_number,
                    current_heading or f"{document_title} Section {item_number}",
                    "\n\n".join(current_chunks),
                    metadata={
                        "kind": row_kind,
                        "heading": current_heading,
                        "heading_level": current_level,
                    },
                )
            )
            current_heading = None
            current_level = None
            current_chunks = []

        body = soup.body or soup
        for element in body.find_all(
            ["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table"]
        ):
            text = element.get_text(" ", strip=True)
            if not text:
                continue

            tag_name = element.name.lower()
            if tag_name.startswith("h") and tag_name[1:].isdigit():
                flush_section()
                current_heading = text
                current_level = int(tag_name[1:])
                current_chunks = [text]
            elif tag_name == "li":
                current_chunks.append(f"- {text}")
            elif tag_name == "table":
                flush_section()
                table_rows = [
                    [cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
                    for row in element.find_all("tr")
                ]
                rows.extend(
                    self._table_to_rows(
                        table_rows,
                        len(rows) + 1,
                        table_index,
                    )
                )
                table_index += 1
            else:
                current_chunks.append(text)

        flush_section()
        if rows:
            return rows

        text = body.get_text(" ", strip=True)
        return [
            DocumentRow(
                1,
                f"{document_title} Section 1",
                text,
                metadata={"kind": row_kind},
            )
        ] if text else []

    def _scrape_txt(self, path):
        text = path.read_text(encoding="utf-8")
        blocks = self._split_blocks(text) or [text.strip()]
        return [
            DocumentRow(index, f"Text Block {index}", block)
            for index, block in enumerate(blocks, start=1)
            if block
        ]

    def _scrape_docx(self, path):
        from docx import Document
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        document = Document(path)
        rows = []
        item_number = 1
        section_title = None
        section_level = None
        section_chunks = []
        paragraph_index = 1
        table_index = 1

        def flush_section():
            nonlocal item_number, section_title, section_level, section_chunks
            if not section_chunks:
                return

            rows.append(
                DocumentRow(
                    item_number,
                    section_title or f"Document Section {item_number}",
                    "\n\n".join(section_chunks),
                    metadata={
                        "kind": "section",
                        "heading_level": section_level,
                    },
                )
            )
            item_number += 1
            section_title = None
            section_level = None
            section_chunks = []

        for child in document.element.body.iterchildren():
            if isinstance(child, CT_P):
                block = Paragraph(child, document)
                text = block.text.strip()
                if not text:
                    continue

                style_name = block.style.name if block.style else ""
                heading_level = self._heading_level(style_name)
                if heading_level is not None:
                    flush_section()
                    section_title = text
                    section_level = heading_level
                    section_chunks = [text]
                    continue

                if section_title:
                    section_chunks.append(text)
                else:
                    rows.append(
                        DocumentRow(
                            item_number,
                            f"Paragraph {paragraph_index}",
                            text,
                            metadata={
                                "kind": "paragraph",
                                "paragraph_index": paragraph_index,
                            },
                        )
                    )
                    item_number += 1
                    paragraph_index += 1
            elif isinstance(child, CT_Tbl):
                flush_section()
                table = Table(child, document)
                table_rows = self._table_to_rows(table.rows, item_number, table_index)
                rows.extend(table_rows)
                item_number += len(table_rows)
                table_index += 1

        flush_section()
        return rows

    def _scrape_pdf(self, path):
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        rows = []
        needs_ocr = False
        for index, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if not page_text:
                needs_ocr = True
            rows.append(
                DocumentRow(
                    index,
                    f"Page {index}",
                    page_text,
                    metadata={
                        "kind": "page",
                        "page_number": index,
                        "extraction": "text-layer" if page_text else "ocr-needed",
                    },
                )
            )

        if needs_ocr:
            ocr_rows = self._perform_ocr_on_pdf(path)
            if ocr_rows:
                merged_rows = []
                for row_index, row in enumerate(rows, start=1):
                    ocr_row = ocr_rows[row_index - 1] if row_index - 1 < len(ocr_rows) else None
                    if row.primary_text.strip():
                        merged_rows.append(row)
                    elif ocr_row and self._has_content(ocr_row):
                        merged_rows.append(
                            DocumentRow(
                                row.item_number,
                                row.title,
                                ocr_row.primary_text,
                                ocr_row.secondary_text,
                                {
                                    "kind": "page",
                                    "page_number": row.item_number,
                                    "extraction": "ocr",
                                },
                            )
                        )
                    else:
                        merged_rows.append(row)
                rows = merged_rows
        return rows

    def _scrape_html(self, path):
        from bs4 import BeautifulSoup

        html = path.read_text(encoding="utf-8")
        soup = BeautifulSoup(html, "html.parser")
        title = soup.title.get_text(strip=True) if soup.title else path.stem
        return self._html_to_rows(soup, title, "html_section")

    def _scrape_rtf(self, path):
        from striprtf.striprtf import rtf_to_text

        rtf_content = path.read_text(encoding="utf-8")
        text = rtf_to_text(rtf_content)
        blocks = self._split_blocks(text) or [text.strip()]
        return [
            DocumentRow(index, f"RTF Block {index}", block)
            for index, block in enumerate(blocks, start=1)
            if block
        ]

    def _scrape_epub(self, path):
        from bs4 import BeautifulSoup
        from ebooklib import ITEM_DOCUMENT, epub

        book = epub.read_epub(str(path))
        rows = []
        item_number = 1
        for item in book.get_items():
            if item.get_type() != ITEM_DOCUMENT:
                continue
            if hasattr(item, "is_chapter") and not item.is_chapter():
                continue
            soup = BeautifulSoup(item.get_content(), "html.parser")
            title = (
                soup.title.get_text(strip=True)
                if soup.title
                else getattr(item, "title", "") or f"Chapter {item_number}"
            )
            body = soup.body or soup
            text_parts = [
                element.get_text(" ", strip=True)
                for element in body.find_all(["h1", "h2", "h3", "p", "li"])
                if element.get_text(" ", strip=True)
            ]
            text = "\n\n".join(text_parts) or soup.get_text(" ", strip=True)
            if not text:
                continue
            rows.append(
                DocumentRow(
                    item_number,
                    title,
                    text,
                    metadata={
                        "kind": "chapter",
                        "chapter_number": item_number,
                        "file_name": item.get_name(),
                    },
                )
            )
            item_number += 1
        return rows

    def _scrape_spreadsheet(self, path):
        import pandas as pd

        if path.suffix.lower() == ".csv":
            sheets = {"Sheet1": pd.read_csv(path)}
        else:
            sheets = pd.read_excel(path, sheet_name=None)

        rows = []
        item_number = 1
        for sheet_name, frame in sheets.items():
            cleaned_frame = frame.fillna("")
            for row_index, row in cleaned_frame.iterrows():
                values = [
                    f"{column}: {str(value).strip()}"
                    for column, value in row.items()
                    if str(value).strip()
                ]
                if not values:
                    continue
                rows.append(
                    DocumentRow(
                        item_number,
                        f"{sheet_name} Row {row_index + 1}",
                        "\n".join(values),
                        f"Sheet: {sheet_name}\nColumns: {', '.join(map(str, cleaned_frame.columns))}",
                        {
                            "kind": "sheet_row",
                            "sheet_name": str(sheet_name),
                            "row_number": int(row_index + 1),
                            "columns": [str(column) for column in cleaned_frame.columns],
                        },
                    )
                )
                item_number += 1
        return rows

    def _scrape_pptx(self, path):
        from pptx import Presentation

        presentation = Presentation(str(path))
        rows = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            notes_text = ""
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            rows.append(
                DocumentRow(
                    index,
                    f"Slide {index}",
                    " ".join(slide_text).strip(),
                    notes_text,
                    {
                        "kind": "slide",
                        "slide_number": index,
                        "has_notes": bool(notes_text),
                    },
                )
            )
        return rows

    def _scrape_image(self, path):
        ocr_text = self._perform_ocr_on_image(path)
        return [
            DocumentRow(
                1,
                path.stem or "Image 1",
                ocr_text,
                metadata={
                    "kind": "image",
                    "extraction": "ocr",
                },
            )
        ]

    def _perform_ocr_on_pdf(self, path):
        try:
            import pypdfium2 as pdfium
        except ImportError as error:
            raise ValueError(
                "OCR support for scanned PDFs requires the 'pypdfium2' package."
            ) from error

        document = pdfium.PdfDocument(str(path))
        rows = []
        for index in range(len(document)):
            page = document[index]
            bitmap = page.render(scale=2).to_pil()
            page_text = self._ocr_image_object(bitmap)
            rows.append(DocumentRow(index + 1, f"Page {index + 1}", page_text))
        return rows

    def _perform_ocr_on_image(self, path):
        try:
            from PIL import Image
        except ImportError as error:
            raise ValueError(
                "OCR support for image files requires the 'Pillow' package."
            ) from error

        with Image.open(path) as image:
            return self._ocr_image_object(image)

    def _ocr_image_object(self, image):
        try:
            import pytesseract
        except ImportError as error:
            raise ValueError(
                "OCR support requires the 'pytesseract' package and a Tesseract installation."
            ) from error

        return pytesseract.image_to_string(image).strip()
