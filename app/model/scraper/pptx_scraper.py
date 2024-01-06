"""pptx_scraper.py

This module contains the PPTXScraper class, which is designed to extract text
and notes from PowerPoint (PPTX) presentations and optionally save this data to
a CSV file.

The PPTXScraper class takes the path of a PowerPoint file as input, reads
through each slide, and collects the text and notes. This can be particularly
useful for extracting information from presentations for analysis, archiving,
or processing.

Example:
    pptx_path = "path_to_your_pptx_file.pptx"
    csv_path = "output_file.csv"
    scraper = PPTXScraper(pptx_path)
    scraper.save_to_csv(csv_path)
"""

import csv
from pptx import Presentation

class PPTXScraper:
    """
    A class for scraping text and notes from PowerPoint (PPTX) files.

    This class provides methods to extract text from each slide and its
    corresponding notes, then compile this information into a structured
    format, such as a list or a CSV file.
    """

    def __init__(self, pptx_path):
        """
        Initialize the PPTXScraper with the path to a PowerPoint file.

        Args:
            pptx_path (str): Path to the PowerPoint (PPTX) file to be processed.
        """
        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)

    def _get_slide_text(self, slide):
        """
        Extract text from a single slide.

        Args:
            slide (Slide): A slide object from the PowerPoint presentation.

        Returns:
            str: All text found in the slide, concatenated and separated by spaces.
        """
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
        return text.strip()

    def _get_slide_notes(self, slide):
        """
        Extract notes from a single slide.

        Args:
            slide (Slide): A slide object from the PowerPoint presentation.

        Returns:
            str: Text from the notes section of the slide, if available.
        """
        if slide.has_notes_slide:
            return slide.notes_slide.notes_text_frame.text
        return ""

    def scrape(self):
        """
        Scrape text and notes from all slides in the presentation.

        Returns:
            list: A list of tuples, each containing the slide number, slide text, and notes text.
        """
        data = []
        for i, slide in enumerate(self.presentation.slides, start=1):
            slide_text = self._get_slide_text(slide)
            notes_text = self._get_slide_notes(slide)
            data.append((i, slide_text, notes_text))
        return data

    def save_to_csv(self, csv_path):
        """
        Save the scraped data to a CSV file.

        Args:
            csv_path (str): Path where the CSV file will be saved.
        """
        data = self.scrape()
        with open(csv_path, "w", newline="", encoding="utf-8") as file:
            writer = csv.writer(file)
            writer.writerow(["Slide Number", "Slide Text", "Notes"])
            writer.writerows(data)
